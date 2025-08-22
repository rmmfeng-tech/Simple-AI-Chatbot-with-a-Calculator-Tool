import os

import sys

import subprocess

import json

import pptx

from pathlib import Path

from pptx.util import Cm

from pptx.dml.color import RGBColor 

from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE

from pptx.enum.text import PP_ALIGN

def handle_tool_call(message):
    name = message[0].function.name
    if name == 'calculator':
        arguments = message[0].function.arguments
        operation,a,b = arguments.values()
        result = calculator(operation,a,b)
        response = {
            "role": "tool",
            "content": json.dumps({"result": result}),
        }
        return response
    elif name == 'word_counter':
        arguments = message[0].function.arguments
        text = arguments['text']
        words = word_counter(text)
        response = {
            "role": "tool",
            "content": json.dumps({"result": words}),
        }
        return response
    elif name == 'characters_count':
        arguments = message[0].function.arguments
        character = arguments['characters']
        result = characters_count(character)
        response = {
            "role": "tool",
            "content": json.dumps({"result": result}),
        }
    elif name == 'simple_fluxogram':
        arguments = message[0].function.arguments
        list_of_tasks = arguments['list_of_tasks']
        if arguments['file_name']:
            file_name = arguments['file_name']
            result = simple_fluxogram(list_of_tasks,file_name)
        else:
            result = simple_fluxogram(list_of_tasks)
        response = {
            "role": "tool",
            "content": json.dumps({"result": result}),
        }
        return response
        
def calculator(operation:str,a:float,b:float):
    """
    # simple calculator tool
    ----------
    ## Paramaters
    ----------
    operation : str
        this paramater defines what calculation will be done, the options are
        - 'sum'
        - 'subtraction'
        - 'multiplication'
        - 'division'
        - 'raising to power'
    
    a : float
        first value of the calculation

    b : float
        second value of the calculation

    ## Example of usage
   ----------
    **sum example**: calculator(operation ='sum',a=2,b=3)
        result = 2+3
        return result
    **subtraction example**: calculator(operation ='subtraction',a=2,b=3)
        result = 2-3
        return result
    **multiplication example**: calculator(operation ='multiplication',a=2,b=3)
        result = 2*3
        return result
    **division example**: calculator(operation ='division',a=2,b=3)
        result = a/b
        return result
    **raising to power example**: calculator(operation ='raising to power',a=2,b=3)
        result = a^b
        return result
    """
    result = 0
    if operation == 'sum':
        result = a+b
    elif operation == 'subtraction':
        result = a-b
    elif operation == 'multiplication':
        result = a*b
    elif operation == 'division':
        result = a/b
    elif operation == 'raising to power':
        result = a^b
    return result

def word_counter(text: str):
    """
    # simple word counter tool
    ----------
    ## Parameters
    ----------
    text : str
        The text whose words you want to count.

    ## Example of usage
    ----------
    word_counter('Hello world') 
        words = 2
        return len(words)
    """
    print(text)
    words = text.split()
    return len(words)

def characters_count(characters):
    """
    # simple characters counter tool
    ----------
    ## Paramaters
    ----------
    characters : str 
        The string whose characters you want to count.

    ## Example of usage
        count_characters('123456') 
        characters = 6
        return len(characters)
    """
    return len(characters)

def simple_fluxogram(list_of_tasks,file_name="fluxogram"):
    """
    # simple fluxogram creator
    ----------
    ## Paramaters
    ----------
    list_of_tasks : list
        list of tasks in order of execution
    file_name : str
        optinal name for the archive, default is fluxogram

    ## Example of usage
        list_of_tasks = [fix the screw, polish the edges]
        simple_fluxogram(list_of_tasks) 
        file = file_name+".pptx"
        template.save(file)
        os.startfile(file)
        os.startfile(os.path.dirname(file))
    """
    division = Cm(3)
    shape_height = Cm(5)
    shape_width = Cm(10)
    left = Cm(5)  
    top = Cm(1)   
    template = pptx.Presentation()
    slide_layout = template.slide_layouts[6]
    number_of_shapes = len(list_of_tasks)
    i = ((shape_height + division) * number_of_shapes) + 2*top - division
    slide = template.slides.add_slide(slide_layout)   
    template.slide_width = Cm(20)
    template.slide_height = i
    for item in list_of_tasks:
        retangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,left,top,shape_width,shape_height)
        text = retangle.text_frame 
        text.text = str(item)
        text.paragraphs[0].alignment = PP_ALIGN.CENTER
        retangle.fill.solid()
        retangle.fill.fore_color.rgb = RGBColor(79,129,189)
        if item !=list_of_tasks[-1]:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Cm(9.25),top+shape_height,division*.5,division)
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(192,80,77)
        top += shape_height + division
    
    file = file_name+".pptx"
    folder = os.path.dirname(file)
    template.save(file)
    if sys.platform == "win32":
        os.startfile(file)
        os.startfile(folder)
    elif sys.platform == "darwin":
        subprocess.run(["open", file])
        subprocess.run(["open", folder])
    else:
        subprocess.run(["xdg-open", file])
        subprocess.run(["xdg-open", folder])

    return 'pptx file created'

simple_fluxogram_function = {
    "name": "simple_fluxogram",
    "description": "Creates a simple fluxogram in a pptx archive (only works on windows).",
    "parameters": {
        "type": "object",
        "properties": {
                        "list_of_tasks" : {"type": "list", "description": "list of tasks in order of execution"},
            "file_name" : {"type": "string", "description": "optinal name for the archive, default is fluxogram"}
        },
        "required": ["list_of_tasks"],
    },
}

characters_count_function = {
    "name": "characters_count",
    "description": "Counts the number of characters in a given text.",
    "parameters": {
        "type": "object",
        "properties": {
            "characters": {"type": "string", "description": "The string whose characters you want to count."}
        },
        "required": ["characters"],
    },
}
# definition of the function in a dict so the AI knows how to use it
word_counter_function = {
    "name": "word_counter",
    "description": "Counts the number of words in a given text.",
    "parameters": {
        "type": "object",
        "properties": {
            "text": {"type": "string", "description": "The text whose words you want to count."}
        },
        "required": ["text"],
    },
}
#definiton of the function in a dict so the AI know how to use it
calculator_function = {
    "name": "calculator",
    "description": """
    Simple calculator function that make the following operations :
     * sum
     * subtraction
     * multiplication
     * division
     * raising to power
    """,
    "parameters": {
        "type": "object",
        "properties": {
            'operation': {'type': 'string','description': 'the tipe of calculation you want to do'},
            'a': {'type': 'float', 'description': 'The first number in the calculation'},
            'b': {'type': 'float', 'description': 'The second number in the calculation'},
        },
        "required": ['operation','a','b'],
    }
}

tools = [
    {"type": "function", "function": calculator_function},
    {"type": "function", "function": word_counter_function},
    {"type": "function", "function": characters_count_function},
    {"type": "function", "function": simple_fluxogram_function}
    ]